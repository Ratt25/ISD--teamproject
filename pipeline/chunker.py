"""
Material 파일 → 텍스트 추출 → Doc_Chunk 저장
Usage: python -m pipeline.chunker
"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from db.db import get_conn, init_schema, insert_doc_chunks, chunks_exist

MIN_CHARS = 30  # 너무 짧은 페이지는 스킵


def extract_pdf(path: Path) -> list[dict]:
    import fitz
    doc = fitz.open(str(path))
    chunks = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if len(text) >= MIN_CHARS:
            chunks.append({"content": text, "page_ref": i + 1, "chunk_index": i})
    doc.close()
    return chunks


def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text.strip())
        text = "\n".join(t for t in texts if t)
        if len(text) >= MIN_CHARS:
            chunks.append({"content": text, "page_ref": i + 1, "chunk_index": i})
    return chunks


def chunk_material(material_id: int, file_path: str, file_type: str) -> int:
    """Material 1개 청킹 → Doc_Chunk 삽입. 반환값: 저장된 청크 수."""
    if chunks_exist(material_id):
        return 0

    path = Path(file_path)
    if not path.exists():
        return 0

    ft = (file_type or "").lower()
    try:
        if ft == "pdf":
            chunks = extract_pdf(path)
        elif ft in ("pptx", "ppt"):
            chunks = extract_pptx(path)
        else:
            return 0
    except Exception as e:
        print(f"  [WARN] 추출 실패 {path.name}: {e}")
        return 0

    if not chunks:
        return 0

    insert_doc_chunks(material_id, chunks)
    return len(chunks)


def run_chunker() -> dict:
    """DB의 모든 Material을 순회해 Doc_Chunk가 없는 것만 청킹."""
    init_schema()
    with get_conn() as conn:
        materials = conn.execute(
            "SELECT material_id, file_path, file_type, title FROM Material"
        ).fetchall()

    total = 0
    skipped = 0
    for m in materials:
        n = chunk_material(m["material_id"], m["file_path"], m["file_type"])
        if n:
            print(f"  [chunk] {Path(m['file_path']).name} → {n}청크")
            total += n
        else:
            skipped += 1

    print(f"\n완료: {total}청크 생성, {skipped}개 스킵")
    return {"chunks": total, "skipped": skipped}


if __name__ == "__main__":
    print(run_chunker())
